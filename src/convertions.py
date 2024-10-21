# Import dependencies ; if not installed yet: pip install python-pptx pdfplumber
from pptx import Presentation
import pdfplumber


def _get_extension (file_path):
    """
        This function takes a file path as an argument and extract its extension
    Args:
        - file_path (String): Exact file path of pdf file that is accourding to source python file.
    Returns:
        - extension (String): File extension.
    """
    extension = (file_path.split('.'))[-1]
    return extension


def _pdf2text (file_path):
    """
        This function takes a file path as an argument, that file is expected a pdf file and convert it to the text and returns it.
    Args:
        - file_path (String): Exact file path of pdf file that is accourding to source python file.
    Returns:
        - extracted_text (String): Text that extracted from pdf file.
    """
    with pdfplumber.open(file_path) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text()
    return text
            

def _pptx2text (file_path):
    """
        This function takes a file path as an argument, that file is expected a pptx file and convert it to the text and returns it. 
    Args: 
        - file_path (String): Exact file path of pdf file that is accourding to source python file.
    Returns:
        - extracted_text (String): Text that extracted from pdf file.
    """
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()


def convert2txt (file_path):
    """
        - This function takes a file path as an argument, that file is expected a pptx or pdf file and convert it to the text and returns it. 
    Args:
        - file_path (String): Exact file path of pdf file that is accourding to source python file.
    Returns:
        - extracted_text (String): Text that extracted from pdf file.
    """
    extension = _get_extension(file_path)
    text = ""
    if extension == "pdf":
        text = _pdf2text(file_path)
    elif extension == "pptx":
        text = _pptx2text(file_path)
    else:
        print("Unsupported file type!")
        
    return text