import argparse  # Library for parsing command-line arguments
import PyPDF2  # Library for reading PDF files
import json  # Library for working with JSON data
from pptx import Presentation  # Library for working with PowerPoint presentations

# Function to extract text from a PDF file
def extract_text_from_pdf(pdf_path):  
    """
    Converts the contents of a PDF file to a string.
    """
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        extracted_text = ""
        # Loop through each page in the PDF and extract text
        for page in pdf_reader.pages:
            text = page.extract_text()
            if text:
                extracted_text += text  # Append text if present
        return extracted_text

# Function to extract text from a PowerPoint (.pptx) file
def extract_text_from_pptx(pptx_path):  
    prs = Presentation(pptx_path)
    text_list = ""

    # Loop through each slide and each shape to extract text
    for slide_number, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame:  # Check if shape has text content
                text_list += shape.text + "\n"  # Append text with a newline
    return text_list

# Set up command-line argument parser with description
parser = argparse.ArgumentParser(description="Extract text from the given file")
parser.add_argument("file_path", type=str)  # Path to the file to be processed
args = parser.parse_args()

# Store the file path and process the file extension
sourceFile = args.file_path
out_path = sourceFile.split(".")
file_type = out_path[-1]  # Identify the file type by its extension
out_path[-1] = "json"  # Set output file extension to JSON
out_path = ".".join(out_path)  # Reconstruct the output file path

text_output = ""  # Initialize the output text

# Check the file type and call the appropriate function
if file_type == "pdf":  
    text_output = extract_text_from_pdf(sourceFile)
if file_type == "pptx":
    text_output = extract_text_from_pptx(sourceFile)

# Create a dictionary with the source file and extracted text
veri = {
    "source": sourceFile,
    "text": text_output,
}

# Write the extracted text to a new JSON file
with open(out_path, 'w') as f:
    json.dump(veri, f, ensure_ascii=False, indent=4)

# Print the path of the output JSON file
print(out_path)
