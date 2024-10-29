from pptx import Presentation
import os

import google.generativeai as genai


pptx_path = "Chap8_02.pptx" #siteye input olarak verşlen dosyanın path'ı otomatik olarak gelmeli
prs = Presentation(pptx_path)

text_list = ""

for slide_number, slide in enumerate(prs.slides):
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_list += shape.text + "\n"

with open("deneme.txt","w") as f:
    f.write(text_list + "\n")
text_path = "deneme.txt"
#sınav formlarımız olmalı, o formlardan hangisi örnek alınarak exam oluşturulacak bunu kullanıcı seçmeli

genai.configure(api_key="AIzaSyCStzQ1Pv37EINES5z9HAybwSQ94sQGKdw")
model = genai.GenerativeModel("gemini-1.5-flash")
response = model.generate_content([text_list,"Create a well-structured exam based on the provided lecture notes.\
The exam should assess students' understanding of the key concepts, their ability to apply these concepts to new situations, and their critical thinking skills. Create an answer key at the and for multi-choice and true-false questions."])

with open("exam.txt","w") as e:
    e.write(response.text)
